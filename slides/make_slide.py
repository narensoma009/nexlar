"""Generate an executive 16:9 slide for QuoteIQ.

Usage:  python make_slide.py
Output: QuoteIQ.pptx (one slide, editable in PowerPoint)
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# ─── palette ──────────────────────────────────────────────────────────────
BG_TOP = RGBColor(0x0F, 0x17, 0x2A)
BG_BOTTOM = RGBColor(0x1E, 0x29, 0x3B)
ACCENT_HERO = RGBColor(0x6E, 0x8A, 0xFF)
TEXT_PRIMARY = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_MUTED = RGBColor(0xCB, 0xD5, 0xE1)
TEXT_DIM = RGBColor(0x94, 0xA3, 0xB8)
RULE = RGBColor(0x33, 0x41, 0x55)

AGENT_LEFT = RGBColor(0x4F, 0x46, 0xE5)
AGENT_RIGHT = RGBColor(0x7C, 0x3A, 0xED)
AGENT_INK = RGBColor(0xFF, 0xFF, 0xFF)
AGENT_TAG = RGBColor(0xC7, 0xD2, 0xFE)

OUTCOME_FILL = RGBColor(0x06, 0x4E, 0x3B)
OUTCOME_STROKE = RGBColor(0x10, 0xB9, 0x81)
OUTCOME_NUM = RGBColor(0x34, 0xD3, 0x99)
OUTCOME_CAP = RGBColor(0x6E, 0xE7, 0xB7)
OUTCOME_SUB = RGBColor(0xD1, 0xFA, 0xE5)


# ─── helpers ──────────────────────────────────────────────────────────────
def add_text(
    slide, left, top, width, height, text,
    size=14, bold=False, color=TEXT_PRIMARY,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI",
    spacing=1.05,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def stroke(shape, color, weight=1.25):
    shape.line.color.rgb = color
    shape.line.width = Pt(weight)


def add_rect(slide, left, top, width, height, fill_color, radius=False,
             line_color=None, line_weight=1.25):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.shadow.inherit = False
    fill(shape, fill_color)
    if line_color is None:
        no_line(shape)
    else:
        stroke(shape, line_color, line_weight)
    return shape


def add_agent_card(slide, x, y, w, h, color, number, title_lines, tagline, bullets):
    card = add_rect(slide, x, y, w, h, color, radius=True)
    card.adjustments[0] = 0.08  # rounder corners
    # Number badge
    add_text(slide, x + Inches(0.22), y + Inches(0.18), Inches(2.5), Inches(0.3),
             number, size=11, bold=True, color=AGENT_TAG)
    # Title (multi-line)
    title_top = y + Inches(0.5)
    for i, line in enumerate(title_lines):
        add_text(slide, x + Inches(0.22), title_top + Inches(0.32 * i),
                 w - Inches(0.4), Inches(0.4), line,
                 size=18, bold=True, color=AGENT_INK)
    # Tagline
    add_text(slide, x + Inches(0.22), y + Inches(1.25),
             w - Inches(0.4), Inches(0.3),
             tagline, size=11, bold=False, color=AGENT_TAG)
    # Divider line
    div = slide.shapes.add_connector(1, x + Inches(0.22), y + Inches(1.58),
                                     x + w - Inches(0.22), y + Inches(1.58))
    div.line.color.rgb = AGENT_TAG
    div.line.width = Pt(0.5)
    # Bullets
    bullets_box = slide.shapes.add_textbox(
        x + Inches(0.22), y + Inches(1.68),
        w - Inches(0.4), h - Inches(1.8),
    )
    tf = bullets_box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        r = p.add_run()
        r.text = f"•  {b}"
        r.font.name = "Segoe UI"
        r.font.size = Pt(11.5)
        r.font.color.rgb = RGBColor(0xE0, 0xE7, 0xFF)


# ─── slide ────────────────────────────────────────────────────────────────
def main() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Background — two stacked rectangles to fake a soft gradient
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height // 2, BG_TOP)
    add_rect(slide, 0, prs.slide_height // 2, prs.slide_width, prs.slide_height // 2, BG_BOTTOM)

    # Decorative accent ring (top-right) — using oval
    accent = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(-2.4),
                                    Inches(4.2), Inches(4.2))
    fill(accent, RGBColor(0x6E, 0x8A, 0xFF))
    accent.fill.transparency = 0  # solid; we set opacity via XML below
    accent.line.fill.background()
    # opacity via XML
    from pptx.oxml.ns import qn
    sppr = accent.fill._xPr.find(qn("a:solidFill"))
    if sppr is not None:
        srgb = sppr.find(qn("a:srgbClr"))
        if srgb is not None:
            alpha = srgb.makeelement(qn("a:alpha"), {"val": "8000"})
            srgb.append(alpha)

    # Logo dot
    logo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.75), Inches(0.75),
                                  Inches(0.55), Inches(0.55))
    fill(logo, ACCENT_HERO)
    no_line(logo)
    add_text(slide, Inches(0.75), Inches(0.75), Inches(0.55), Inches(0.55),
             "Q", size=22, bold=True, color=TEXT_PRIMARY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Eyebrow + sub-eyebrow
    add_text(slide, Inches(1.45), Inches(0.78), Inches(6), Inches(0.3),
             "CAPSTONE · AGENTIC AI", size=11, bold=True, color=TEXT_DIM)
    add_text(slide, Inches(1.45), Inches(1.04), Inches(6), Inches(0.3),
             "QuoteIQ", size=14, bold=False, color=TEXT_MUTED)

    # Hero headline
    add_text(slide, Inches(0.75), Inches(1.75), Inches(12), Inches(0.9),
             "Quotes that approve themselves",
             size=46, bold=True, color=TEXT_PRIMARY)
    add_text(slide, Inches(0.75), Inches(2.7), Inches(12), Inches(0.5),
             "Four cooperating agents collapse the quote-to-approval cycle —",
             size=18, bold=False, color=TEXT_MUTED)
    add_text(slide, Inches(0.75), Inches(3.02), Inches(12), Inches(0.5),
             "eradicating re-keying, ASC-606 surprises, and the 1.8× resubmit loop.",
             size=18, bold=False, color=TEXT_MUTED)

    # Thin separator
    sep = slide.shapes.add_connector(1, Inches(0.75), Inches(3.6),
                                     Inches(12.58), Inches(3.6))
    sep.line.color.rgb = RULE
    sep.line.width = Pt(0.75)

    # Section eyebrow
    add_text(slide, Inches(0.75), Inches(3.74), Inches(6), Inches(0.3),
             "THE INTELLIGENT LAYER", size=11, bold=True, color=TEXT_DIM)

    # ─── Agent cards (4) ─────────────────────────────────────────────────
    card_w = Inches(2.95)
    card_h = Inches(2.6)
    card_top = Inches(4.05)
    gap = Inches(0.18)
    left0 = Inches(0.75)

    agents = [
        (AGENT_LEFT, "1   COMPOSE",
         ["Catalogue-driven", "composition"],
         "Eradicates Excel re-keying",
         ["Semantic SKU search",
          "Phase-aware suggestions",
          "Inline justify, no re-keying",
          "Validated building blocks"]),
        (AGENT_LEFT, "2   VALIDATE",
         ["Upstream rule", "enforcement"],
         "Catches breaks at config time",
         ["Phasing & sequencing",
          "ASC-606 compliance",
          "DHI codes in plain English",
          "Delta re-check on changes"]),
        (AGENT_RIGHT, "3   DECIDE",
         ["Risk-scored", "auto-routing"],
         "Reserves humans for judgment",
         ["Transparent 0–100 score",
          "Threshold + blockers + signal",
          "Auto-approve in-policy deals",
          "Escalates only the rest"]),
        (AGENT_RIGHT, "4   BRIEF",
         ["LLM approver", "brief"],
         "Decision under 60 seconds",
         ["Pros vs. cons inline",
          "2-sentence recommendation",
          "Mitigations for the AE",
          "Grounded in your rules"]),
    ]
    for i, (color, num, title_lines, tag, bullets) in enumerate(agents):
        left = left0 + (card_w + gap) * i
        add_agent_card(slide, left, card_top, card_w, card_h,
                       color, num, title_lines, tag, bullets)

    # ─── Outcomes ribbon ─────────────────────────────────────────────────
    ribbon = add_rect(slide, Inches(0.75), Inches(6.85),
                      Inches(11.83), Inches(0.5),
                      OUTCOME_FILL, radius=True,
                      line_color=OUTCOME_STROKE, line_weight=1.25)
    ribbon.adjustments[0] = 0.3

    # Eyebrow
    add_text(slide, Inches(0.95), Inches(6.88), Inches(2), Inches(0.3),
             "OUTCOMES", size=10, bold=True, color=OUTCOME_CAP)

    # Outcome cells
    cells = [
        ("↓ days",    "Cycle time"),
        ("≈ 0",       "Reject & resubmit loop"),
        ("100%",      "ASC-606 at submission"),
        ("Auto",      "Approve clean & in-policy"),
    ]
    base_left = Inches(2.55)
    cell_w = Inches(2.5)
    for i, (big, small) in enumerate(cells):
        left = base_left + cell_w * i
        add_text(slide, left, Inches(6.88), Inches(2), Inches(0.3),
                 big, size=14, bold=True, color=OUTCOME_NUM,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, left + Inches(0.9), Inches(6.88), cell_w - Inches(0.9), Inches(0.3),
                 small, size=10, bold=False, color=OUTCOME_SUB,
                 anchor=MSO_ANCHOR.MIDDLE)

    # ─── Footer ──────────────────────────────────────────────────────────
    add_text(slide, Inches(0.75), Inches(7.18), Inches(8), Inches(0.25),
             "QuoteIQ · Pod 6 Capstone · powered by Azure OpenAI Foundry",
             size=10, bold=False, color=TEXT_DIM)
    add_text(slide, Inches(11.6), Inches(7.18), Inches(1), Inches(0.25),
             "1 / 1", size=10, bold=False, color=TEXT_DIM, align=PP_ALIGN.RIGHT)

    out = Path(__file__).parent / "QuoteIQ.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    out = main()
    print(f"Wrote {out}")
